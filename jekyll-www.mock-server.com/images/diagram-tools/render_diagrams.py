#!/usr/bin/env python3
"""Render the generated diagram slides of MockServerScenarios.pptx to PNGs in the
MockServer house style (dashed stacked Expectation boxes, white shadowed boxes,
coloured arrows with arrowheads, Helvetica text), driven by the real pptx
geometry. Read-only on the .pptx; writes PNGs into the images dir.

Why not export from PowerPoint? On a sandboxed macOS PowerPoint, scripted/GUI
export silently drops file writes, so this renders the vector geometry directly.
For pixel-perfect FINALS, open the .pptx in PowerPoint and File > Export the
slides yourself (see make_cropped_decks.py and README.md).

To add a diagram: generate its slide with build_diagrams.py, then add an entry to
SLIDES below (0-based slide index -> primary arrow colour + output filename).

Usage:  python3 render_diagrams.py
"""
import os, math
from pathlib import Path
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont, ImageChops

HERE = Path(__file__).resolve().parent
PPTX = HERE.parent / "MockServerScenarios.pptx"
IMAGES = HERE.parent
NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
W = 1600          # output width in px
SS = 2            # supersample factor for crisp anti-aliasing

THEME = {
    'accent1': (0x4F,0x81,0xBD), 'accent4': (0x80,0x64,0xA2),
    'accent5': (0x4B,0xAC,0xC6), 'red': (0xC0,0x00,0x00),
    'black': (0x40,0x40,0x40),
}
FONTS = ["/System/Library/Fonts/Helvetica.ttc",
         "/System/Library/Fonts/Supplemental/Arial.ttf"]

# 0-based slide index -> (primary arrow colour, output filename stem)
# NOTE: this renderer draws the FLOW-style diagrams (Expectation/Proxy boxes,
# matcher, action, cylinders). The "verify responses" diagram is the
# architecture style (Test / System Under Test / Mock Server / Service boxes);
# it is produced by make_response_verification.py instead, not here.
SLIDES = {
    11: ('accent1', 'mockserver_chaos_action'),
    12: ('accent4', 'mockserver_llm_mocking'),
    14: ('accent5', 'mockserver_llm_optimisation'),
}

def font(px):
    for f in FONTS:
        if os.path.exists(f):
            try: return ImageFont.truetype(f, px)
            except Exception: pass
    return ImageFont.load_default()

def resolve_color(col, primary):
    if col is None: return THEME[primary]
    if col.startswith('srgb:'):
        v = col[5:]; return tuple(int(v[i:i+2], 16) for i in (0, 2, 4))
    if col.startswith('scheme:'):
        s = col[7:]
        return THEME['black'] if s == 'tx1' else THEME.get(s, THEME[primary])
    return THEME[primary]

def shape_color(sh):
    for ln in sh._element.iter('{%s}ln' % NS['a']):
        for c in ln.iter('{%s}srgbClr' % NS['a']): return 'srgb:' + c.get('val')
        for c in ln.iter('{%s}schemeClr' % NS['a']): return 'scheme:' + c.get('val')
    return None

def is_dashed(sh):
    return sh._element.find('.//{%s}prstDash' % NS['a']) is not None

def arrow_ends(sh):
    r = {}
    for t in ('headEnd', 'tailEnd'):
        e = sh._element.find('.//{%s}%s' % (NS['a'], t))
        if e is not None and e.get('type') not in (None, 'none'): r[t] = e.get('type')
    return r

def geom_of(sh):
    pg = sh._element.find('.//{%s}prstGeom' % NS['a'])
    return pg.get('prst') if pg is not None else None

def rounded(d, box, radius, outline, width, dash=None):
    r = radius; x0, y0, x1, y1 = box
    if dash is None:
        d.rounded_rectangle(box, radius=r, outline=outline, width=width); return
    def dash_line(p, q):
        ax, ay = p; bx, by = q; L = math.hypot(bx - ax, by - ay)
        if L == 0: return
        for i in range(int(L // (dash * 2)) + 1):
            a = i * dash * 2; b = min(a + dash, L)
            if a >= L: break
            d.line([(ax + (bx - ax) * a / L, ay + (by - ay) * a / L),
                    (ax + (bx - ax) * b / L, ay + (by - ay) * b / L)], fill=outline, width=width)
    dash_line((x0 + r, y0), (x1 - r, y0)); dash_line((x0 + r, y1), (x1 - r, y1))
    dash_line((x0, y0 + r), (x0, y1 - r)); dash_line((x1, y0 + r), (x1, y1 - r))
    for cx, cy, a0, a1 in [(x0+r,y0+r,180,270),(x1-r,y0+r,270,360),(x1-r,y1-r,0,90),(x0+r,y1-r,90,180)]:
        d.arc([cx-r, cy-r, cx+r, cy+r], a0, a1, fill=outline, width=width)

def draw_centered(d, box, text, fnt, color):
    x0, y0, x1, y1 = box; lines = text.split('\n')
    lh = max(d.textbbox((0,0), l, font=fnt)[3] - d.textbbox((0,0), l, font=fnt)[1] for l in lines) + 6*SS
    cy = (y0 + y1) / 2 - lh * len(lines) / 2
    for i, l in enumerate(lines):
        bb = d.textbbox((0,0), l, font=fnt); w = bb[2] - bb[0]
        d.text(((x0+x1)/2 - w/2, cy + i*lh - bb[1]), l, font=fnt, fill=color)

def box_with_shadow(d, box, text, fnt, radius=0):
    x0, y0, x1, y1 = box; so = 7*SS
    d.rectangle([x0+so, y0+so, x1+so, y1+so], fill=(200,200,200))
    if radius > 0: d.rounded_rectangle(box, radius=radius, fill='white', outline=(0,0,0), width=3*SS)
    else: d.rectangle(box, fill='white', outline=(0,0,0), width=3*SS)
    draw_centered(d, box, text, fnt, (0,0,0))

def cylinder(d, box, text, fnt):
    x0, y0, x1, y1 = box; eh = (y1 - y0) * 0.16
    d.rectangle([x0, y0+eh/2, x1, y1-eh/2], fill='white')
    d.line([x0, y0+eh/2, x0, y1-eh/2], fill=(0,0,0), width=3*SS)
    d.line([x1, y0+eh/2, x1, y1-eh/2], fill=(0,0,0), width=3*SS)
    d.arc([x0, y1-eh, x1, y1], 0, 180, fill=(0,0,0), width=3*SS)
    d.ellipse([x0, y0, x1, y0+eh], fill='white', outline=(0,0,0), width=3*SS)
    draw_centered(d, (x0, y0+eh, x1, y1), text, fnt, (0,0,0))

def draw_label(d, box, text, fnt):
    x0, y0, x1, y1 = box; bb = d.textbbox((0,0), text, font=fnt)
    d.text((x0, (y0+y1)/2 - (bb[3]-bb[1])/2 - bb[1]), text, font=fnt, fill=THEME['black'])

def arrowhead(d, tip, frm, color, size):
    ang = math.atan2(tip[1]-frm[1], tip[0]-frm[0]); spread = math.radians(24)
    p1 = (tip[0]-size*math.cos(ang-spread), tip[1]-size*math.sin(ang-spread))
    p2 = (tip[0]-size*math.cos(ang+spread), tip[1]-size*math.sin(ang+spread))
    d.polygon([tip, p1, p2], fill=color)

def render(prs, idx, primary, name):
    s = prs.slides[idx]
    boxes = [sh for sh in s.shapes if sh.left is not None]
    L = min(b.left for b in boxes); T = min(b.top for b in boxes)
    R = max(b.left+b.width for b in boxes); B = max(b.top+b.height for b in boxes)
    GW, GH = R-L, B-T
    cw = W*SS; ch = int(cw*GH/GW)
    img = Image.new('RGB', (cw, ch), 'white'); d = ImageDraw.Draw(img)
    NX = lambda v: (v-L)/GW*cw; NY = lambda v: (v-T)/GH*ch
    fbox = font(int(0.017*cw)); flbl = font(int(0.015*cw)); fexp = font(int(0.020*cw))
    shapes, lines = [], []
    for sh in boxes:
        bx = (NX(sh.left), NY(sh.top), NX(sh.left+sh.width), NY(sh.top+sh.height))
        (lines if sh.shape_type == 9 else shapes).append(
            (sh, geom_of(sh), bx, sh.text_frame.text if sh.has_text_frame else ''))
    matcher_px = None; ret_label_y = None
    for sh, g, bx, txt in shapes:
        if 'Matcher' in txt: matcher_px = bx
        if txt.strip().startswith('5.'): ret_label_y = (bx[1]+bx[3])/2
    for sh, g, bx, txt in shapes:
        if g == 'roundRect' and sh.shape_type == 6:          # dashed stacked Expectation group
            x0, y0, x1, y1 = bx; o = 0.02*cw
            for k in (2, 1, 0):
                rounded(d, (x0+k*o, y0+k*o, x1-(2-k)*o, y1-(2-k)*o),
                        0.06*(y1-y0), (0,0,0), 3*SS, dash=14*SS)
            draw_centered(d, (x0, y0+(y1-y0)*0.62, x1-2*o, y1-2*o), 'Expectation', fexp, (0,0,0))
        elif g == 'roundRect' and is_dashed(sh):             # dashed grey Expectation/Proxy box
            grey = (130,130,130)
            rounded(d, bx, 0.12*(bx[3]-bx[1]), grey, 3*SS, dash=14*SS)
            draw_centered(d, bx, txt.replace('/','\n'), fbox, grey)
        elif g == 'roundRect':
            box_with_shadow(d, bx, txt.replace('/','\n'), fbox, radius=0.10*(bx[3]-bx[1]))
        elif g == 'can':
            cylinder(d, bx, txt.replace('/','\n'), fbox)
        elif g == 'rect' and txt and sh.shape_type == 1:
            box_with_shadow(d, bx, txt.replace('/','\n'), fbox)
        elif txt:
            draw_label(d, bx, txt, flbl)
    for sh, g, bx, _ in lines:
        col = resolve_color(shape_color(sh), primary)
        ends = arrow_ends(sh); x0, y0, x1, y1 = bx; lw = 4*SS; ah = 16*SS
        xf = sh._element.find('.//{%s}xfrm' % NS['a'])
        flipH = xf is not None and xf.get('flipH') == '1'
        flipV = xf is not None and xf.get('flipV') == '1'
        sx, ex = (x1, x0) if flipH else (x0, x1)
        sy, ey = (y1, y0) if flipV else (y0, y1)
        if g == 'bentConnector2' and y1 > 0.9*ch and matcher_px and ret_label_y is not None:
            mcx = (matcher_px[0]+matcher_px[2])/2; ry = ret_label_y + 0.055*ch
            d.line([(mcx, matcher_px[3]), (mcx, ry)], fill=col, width=lw)
            d.line([(mcx, ry), (0, ry)], fill=col, width=lw)
            arrowhead(d, (0, ry), (mcx, ry), col, ah)
        elif g == 'bentConnector2':
            mid = (ex, sy)
            d.line([(sx, sy), mid], fill=col, width=lw); d.line([mid, (ex, ey)], fill=col, width=lw)
            if 'tailEnd' in ends: arrowhead(d, (ex, ey), mid, col, ah)
            if 'headEnd' in ends: arrowhead(d, (sx, sy), mid, col, ah)
        else:
            d.line([(sx, sy), (ex, ey)], fill=col, width=lw)
            if 'tailEnd' in ends: arrowhead(d, (ex, ey), (sx, sy), col, ah)
            if 'headEnd' in ends: arrowhead(d, (sx, sy), (ex, ey), col, ah)
    img = img.resize((W, int(ch/SS)), Image.LANCZOS)
    bbox = ImageChops.difference(img, Image.new('RGB', img.size, 'white')).getbbox()
    if bbox:
        pad = 36; l, t, r, b = bbox
        img = img.crop((max(0,l-pad), max(0,t-pad), min(img.width,r+pad), min(img.height,b+pad)))
    dst = IMAGES / (name + '.png')
    img.save(dst)
    print('rendered', dst.name, img.size)

def main():
    prs = Presentation(str(PPTX))
    for idx, (primary, name) in SLIDES.items():
        render(prs, idx, primary, name)

if __name__ == '__main__':
    main()
