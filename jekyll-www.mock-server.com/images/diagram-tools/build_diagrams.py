#!/usr/bin/env python3
"""Regenerate the NEW diagram slides in MockServerScenarios.pptx by cloning the
original diagram slides, relabelling their text and (optionally) recolouring via
the theme style matrix. Cloning carries every theme/style reference verbatim, so
the new slides inherit the exact house look (dashed Expectation boxes, shadowed
white action boxes, arrowed connectors, Helvetica text).

The first 11 slides are the ORIGINALS and are never touched. This script is
idempotent: it removes any previously-generated slides (12+) and re-adds them
from the definitions below, so re-running always produces the same deck.

WARNING: this overwrites MockServerScenarios.pptx in place. If you hand-tune a
generated slide in PowerPoint for a pixel-perfect export, do NOT re-run this
afterwards or those tweaks are lost — change the definitions here instead.

To add a diagram: append a block in MAIN below (clone a source slide index,
relabel, optionally recolour), give it a slide index, then register that index
in render_diagrams.py SLIDES so it renders to PNG.

Usage:  python3 build_diagrams.py
"""
import copy
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

PPTX = Path(__file__).resolve().parent.parent / "MockServerScenarios.pptx"
ORIGINAL_SLIDES = 11  # slides 1..11 are the originals; everything after is generated
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}
def q(tag):
    pre, name = tag.split(':')
    return '{%s}%s' % (NS[pre], name)

def duplicate_slide(prs, index):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):              # strip layout placeholders
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:                  # deep-copy every source shape
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    return dest

def norm(t):
    return t.replace('\n', '/').strip()

def iter_textframes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:                 # GROUP
            yield from iter_textframes(sh.shapes)
        elif sh.has_text_frame:
            yield sh

def set_text(shape, new_lines):
    """Replace a text frame's contents, preserving the first run's formatting
    (font/size/colour) in each paragraph."""
    txBody = shape.text_frame._txBody
    paras = txBody.findall(q('a:p'))
    while len(paras) < len(new_lines):         # clone last para if we need more
        txBody.append(copy.deepcopy(paras[-1]))
        paras = txBody.findall(q('a:p'))
    for i, line in enumerate(new_lines):
        runs = paras[i].findall(q('a:r'))
        if runs:
            runs[0].find(q('a:t')).text = line
            for r in runs[1:]:
                r.find(q('a:t')).text = ''
    for j in range(len(paras) - 1, len(new_lines) - 1, -1):
        txBody.remove(paras[j])

def relabel(slide, mapping):
    """mapping: { current_text(as 'a/b/c') : [new, lines] }"""
    done = set()
    for sh in iter_textframes(slide.shapes):
        key = norm(sh.text_frame.text)
        if key in mapping and key not in done:
            set_text(sh, mapping[key]); done.add(key)
    missing = set(mapping) - done
    if missing:
        print("   WARNING unmatched labels:", missing)

def recolor(slide, target_accent):
    """Swap accent1 -> target_accent in every colour ref on the slide, keeping
    the themed gradient/shadow/font intact (used to tint a whole diagram)."""
    n = 0
    for el in slide.shapes._spTree.iter(q('a:schemeClr')):
        if el.get('val') == 'accent1':
            el.set('val', target_accent); n += 1
    print(f"   recoloured {n} accent1 refs -> {target_accent}")

def main():
    prs = Presentation(str(PPTX))
    # reset: drop any previously-generated slides (parts + rels + list entries)
    # so this is idempotent and never leaves orphaned/duplicate slide parts
    lst = prs.slides._sldIdLst
    for sldId in list(lst)[ORIGINAL_SLIDES:]:
        prs.part.drop_rel(sldId.get(qn('r:id')))
        lst.remove(sldId)

    # ---- Diagram 1: Chaos / fault injection  (clone slide4 = error action) ----
    print("Chaos / fault injection")
    s = duplicate_slide(prs, 3)
    relabel(s, {'Error/Action': ['Chaos', 'Action'],
                '3. Return Response': ['3. Inject Fault']})

    # ---- Diagram 2: LLM mocking  (clone slide3 = response action) ----
    print("LLM mocking")
    s = duplicate_slide(prs, 2)
    relabel(s, {'Response/Action': ['LLM', 'Response'],
                '1. Receive Request': ['1. Receive Prompt'],
                '3. Return Response': ['3. Stream Tokens']})
    recolor(s, 'accent4')  # purple

    # ---- Diagram 3: Verify responses  (clone slide2 = verify-requests
    # architecture diagram, so it matches the "Verifying Requests" style) ----
    print("Verify responses")
    s = duplicate_slide(prs, 1)
    relabel(s, {'2. Verify Requests': ['2. Verify Responses']})

    # ---- Diagram 4: LLM proxy -> optimise  (clone slide7 = verification) ----
    print("LLM proxy -> optimise")
    s = duplicate_slide(prs, 6)
    relabel(s, {'Expectation/or/Proxy': ['Proxy'],
                'Recorded/Requests': ['LLM', 'Traffic'],
                '1. Receive Request': ['1. Proxy LLM Call'],
                '2. Record Request': ['2. Record Traffic'],
                '3. Send Verification': ['3. Request Brief'],
                '4. Match Request(s)': ['4. Analyse Traffic'],
                '5. Return Result': ['5. Export Brief']})
    recolor(s, 'accent5')  # teal

    prs.save(str(PPTX))
    print(f"\nSaved {PPTX.name} with {len(list(prs.slides._sldIdLst))} slides "
          f"({ORIGINAL_SLIDES} original + generated).")

if __name__ == '__main__':
    main()
